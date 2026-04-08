"""
Dental figure extractor - Vercel Python function.

Receives a PDF binary via POST body, runs PyMuPDF caption-first figure
extraction, and returns a PowerPoint (.pptx) binary with one slide per figure
(image + caption).

Adapted from /Users/anderschan/dental-ai/extract_images.py - same core
caption-first matching logic, but stream-based (no disk I/O) and outputs
pptx instead of ReportLab PDF.
"""

from http.server import BaseHTTPRequestHandler
import io
import json
import re

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt


def _x_overlap(rect1, rect2) -> float:
    r1_x0 = rect1.x0 if hasattr(rect1, "x0") else rect1[0]
    r1_x1 = rect1.x1 if hasattr(rect1, "x1") else rect1[2]
    if isinstance(rect2, dict) and "bbox" in rect2:
        r2_x0 = rect2["bbox"][0]
        r2_x1 = rect2["bbox"][2]
    elif hasattr(rect2, "x0"):
        r2_x0 = rect2.x0
        r2_x1 = rect2.x1
    else:
        return 0
    return max(0, min(r1_x1, r2_x1) - max(r1_x0, r2_x0))


def extract_figures(pdf_bytes: bytes):
    """Caption-first figure extraction. Returns a list of
    {image_bytes: bytes, caption: str, page: int} dicts.
    """
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    figures = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        blocks = page.get_text("dict")["blocks"]

        caption_blocks = []
        for block in blocks:
            if block.get("type") != 0:
                continue
            raw_lines = []
            for line in block.get("lines", []):
                line_raw = ""
                for span in line.get("spans", []):
                    line_raw += span.get("text", "")
                raw_lines.append(line_raw)
            raw_text = " ".join(raw_lines).strip()
            if re.match(r"[•\s]*Fig\.?\s*\d+\.\d+", raw_text):
                caption_blocks.append({
                    "bbox": block["bbox"],
                    "text": raw_text,
                    "y": block["bbox"][1],
                })

        # Collect images on the page
        page_images = []
        for info in page.get_image_info(xrefs=True):
            if not info.get("xref"):
                continue
            img_rect = fitz.Rect(info["bbox"])
            if img_rect.width < 50 or img_rect.height < 50:
                continue
            try:
                pix = page.get_pixmap(clip=img_rect, dpi=200)
                if pix.width < 50 or pix.height < 50:
                    continue
                page_images.append({
                    "rect": img_rect,
                    "bytes": pix.tobytes("png"),
                    "bottom": img_rect.y1,
                    "top": img_rect.y0,
                })
            except Exception:
                continue

        caption_blocks.sort(key=lambda c: c["y"])
        page_images.sort(key=lambda i: i["bottom"])

        # For each caption, find the closest image directly above it
        for cap in caption_blocks:
            cap_top = cap["y"]
            best_img = None
            best_dist = float("inf")
            for img in page_images:
                if img["bottom"] > cap_top + 15:
                    continue
                if _x_overlap(img["rect"], cap) < 10:
                    continue
                dist = cap_top - img["bottom"]
                if 0 <= dist < best_dist:
                    best_dist = dist
                    best_img = img
            if not best_img:
                continue

            fig_images = [best_img]
            page_images.remove(best_img)

            caption_text = re.sub(r"\s+", " ", cap["text"])
            is_multipart = bool(
                re.search(r"\bA\b", caption_text) and re.search(r"\bB\b", caption_text)
            )
            max_gap = 500 if is_multipart else 60

            # Walk up for stacked images above this figure
            curr = best_img
            while True:
                next_img = None
                next_dist = float("inf")
                for img in page_images:
                    if img["bottom"] > curr["top"] + 5:
                        continue
                    if _x_overlap(img["rect"], curr["rect"]) < 10:
                        continue
                    dist = curr["top"] - img["bottom"]
                    if -5 <= dist < max_gap and dist < next_dist:
                        next_dist = dist
                        next_img = img
                if not next_img:
                    break
                fig_images.insert(0, next_img)
                page_images.remove(next_img)
                curr = next_img

            # Sibling detection for multi-part figures side-by-side
            if is_multipart:
                group_top = min(i["top"] for i in fig_images)
                group_bottom = max(i["bottom"] for i in fig_images)
                siblings = []
                for img in page_images:
                    overlap_y = max(
                        0, min(group_bottom, img["bottom"]) - max(group_top, img["top"])
                    )
                    h = img["bottom"] - img["top"]
                    if h > 0 and (overlap_y / h) > 0.5:
                        siblings.append(img)
                for s in siblings:
                    fig_images.append(s)
                    page_images.remove(s)
                fig_images.sort(key=lambda i: (i["top"], i["rect"].x0))

            # De-hyphenate line-wrapped words
            caption_text = re.sub(r"([a-zA-Z])-\s+([a-zA-Z])", r"\1\2", caption_text)

            # Emit one figure entry per image collected (so multi-part figs
            # become multiple slides, each showing the image the caption refers to)
            for idx, img in enumerate(fig_images):
                figures.append({
                    "image_bytes": img["bytes"],
                    "caption": caption_text,
                    "page": page_num + 1,
                    "part": idx + 1 if len(fig_images) > 1 else None,
                })

    doc.close()

    # Sort by Fig number in caption so output follows textbook order
    def sort_key(fig):
        m = re.search(r"Fig\.?\s*(\d+)\.(\d+)", fig["caption"])
        if m:
            return (int(m.group(1)), int(m.group(2)), fig.get("part") or 0)
        return (999, 999, 0)

    figures.sort(key=sort_key)
    return figures


def build_pptx(figures) -> bytes:
    """Generate a 16:9 PowerPoint with one slide per figure.

    Layout: large image on the left (~60% of slide), caption textbox on the right.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    IMG_LEFT = Inches(0.4)
    IMG_TOP = Inches(0.4)
    IMG_MAX_W = Inches(8.2)
    IMG_MAX_H = Inches(6.7)

    CAP_LEFT = Inches(8.9)
    CAP_TOP = Inches(0.5)
    CAP_WIDTH = Inches(4.1)
    CAP_HEIGHT = Inches(6.5)

    for fig in figures:
        slide = prs.slides.add_slide(blank)

        # Figure out native aspect to fit within bounding box
        img_stream = io.BytesIO(fig["image_bytes"])
        try:
            tmp_pic = slide.shapes.add_picture(
                img_stream, IMG_LEFT, IMG_TOP, width=IMG_MAX_W
            )
            if tmp_pic.height > IMG_MAX_H:
                scale = IMG_MAX_H / tmp_pic.height
                tmp_pic.width = int(tmp_pic.width * scale)
                tmp_pic.height = int(tmp_pic.height * scale)
        except Exception:
            # Skip slides where image decode fails
            continue

        tb = slide.shapes.add_textbox(CAP_LEFT, CAP_TOP, CAP_WIDTH, CAP_HEIGHT)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = re.sub(r"<[^>]+>", "", fig["caption"])
        for run in p.runs:
            run.font.size = Pt(12)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", "0"))
            if length <= 0:
                self._json(400, {"error": "empty body"})
                return
            pdf_bytes = self.rfile.read(length)
            figures = extract_figures(pdf_bytes)
            if not figures:
                self._json(422, {"error": "no figures found in PDF"})
                return
            pptx_bytes = build_pptx(figures)
            self.send_response(200)
            self.send_header(
                "Content-Type",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self.send_header("X-Figure-Count", str(len(figures)))
            self.end_headers()
            self.wfile.write(pptx_bytes)
        except Exception as e:
            self._json(500, {"error": str(e)})

    def _json(self, status: int, payload: dict):
        body = json.dumps(payload).encode()
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)
